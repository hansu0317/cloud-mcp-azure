"""
Quali CRM Chat — 고객용 발표자료 생성
python3 scripts/build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x11, 0x17)   # 배경
SURFACE    = RGBColor(0x1A, 0x1D, 0x27)   # 카드
ACCENT     = RGBColor(0x63, 0x66, 0xF1)   # 인디고
ACCENT2    = RGBColor(0x81, 0x8C, 0xF8)   # 연 인디고
GREEN      = RGBColor(0x10, 0xB9, 0x81)   # 초록
YELLOW     = RGBColor(0xF5, 0x9E, 0x0B)   # 노랑
RED        = RGBColor(0xEF, 0x44, 0x44)   # 빨강
TEXT       = RGBColor(0xE2, 0xE8, 0xF0)   # 본문
TEXT_SUB   = RGBColor(0x94, 0xA3, 0xB8)   # 서브텍스트
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)   # 와이드 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # 빈 레이아웃

# ── 헬퍼 함수 ──────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill=SURFACE, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def text(slide, content, x, y, w, h,
         size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
         wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def multiline(slide, lines, x, y, w, h, size=16, color=TEXT, spacing=1.2):
    """lines: list of (text, bold, color)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            content, bold, clr = item, False, color
        else:
            content, bold, clr = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = clr
    return txb

def divider(slide, y, color=ACCENT, x=0.4, w=12.5):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def step_box(slide, num, title, desc, x, y, w=2.8):
    box(slide, x, y, w, 1.5, fill=SURFACE, line_color=ACCENT, line_width=Pt(1.5))
    # 번호 원
    circle = slide.shapes.add_shape(9, Inches(x+0.1), Inches(y+0.1), Inches(0.4), Inches(0.4))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(num); run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = WHITE
    text(slide, title, x+0.6, y+0.08, w-0.7, 0.35, size=13, bold=True, color=ACCENT2)
    text(slide, desc,  x+0.1, y+0.55, w-0.2, 0.85, size=11, color=TEXT_SUB)

def arrow(slide, x, y, horizontal=True):
    if horizontal:
        shape = slide.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.35), Inches(0.25))
    else:
        shape = slide.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.25), Inches(0.35))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT2
    shape.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)

# 좌측 강조 바
box(sl, 0, 0, 0.12, 7.5, fill=ACCENT)

# 배경 장식 원
for cx, cy, cr, alpha in [(10.5,1.2,2.5,0x18), (11.5,5.5,1.8,0x12)]:
    c = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(cr), Inches(cr))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x63,0x66,0xF1)
    c.line.fill.background()

text(sl, "Quali CRM", 0.5, 1.6, 9, 1.0, size=42, bold=True, color=ACCENT2)
text(sl, "AI Chat Assistant", 0.5, 2.5, 9, 0.9, size=38, bold=True, color=WHITE)
divider(sl, 3.55, color=ACCENT, x=0.5, w=6)
text(sl, "Claude Code × MCP × Azure Dataverse", 0.5, 3.7, 10, 0.6, size=20, color=TEXT_SUB)
text(sl, "자연어 질문으로 CRM 데이터를 실시간 조회하는 AI 어시스턴트", 0.5, 4.35, 11, 0.6, size=16, color=TEXT_SUB)
text(sl, "2025", 0.5, 6.7, 3, 0.4, size=13, color=RGBColor(0x4A,0x55,0x68))

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 왜 필요한가 (도입 배경)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "도입 배경", 0.5, 0.3, 8, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

# 기존 방식
box(sl, 0.4, 1.2, 5.6, 5.4, fill=RGBColor(0x1A,0x0A,0x0A), line_color=RED, line_width=Pt(1.5))
text(sl, "❌  기존 방식", 0.7, 1.35, 4, 0.5, size=16, bold=True, color=RED)
multiline(sl, [
    ("• CRM 담당자에게 데이터 요청", False, TEXT_SUB),
    ("• 담당자가 SQL 쿼리 직접 작성", False, TEXT_SUB),
    ("• 결과 Excel 정리 후 공유", False, TEXT_SUB),
    ("• 요청~수신: 수 시간 ~ 하루", False, TEXT_SUB),
    ("", False, TEXT_SUB),
    ("→ 빠른 의사결정 불가", True, RED),
    ("→ 담당자 리소스 낭비", True, RED),
], 0.7, 1.9, 5.0, 4.0, size=15)

# 새로운 방식
box(sl, 6.9, 1.2, 5.9, 5.4, fill=RGBColor(0x05,0x14,0x10), line_color=GREEN, line_width=Pt(1.5))
text(sl, "✅  새로운 방식 (이 시스템)", 7.1, 1.35, 5.2, 0.5, size=16, bold=True, color=GREEN)
multiline(sl, [
    ("• 웹 채팅창에 자연어로 질문", False, TEXT_SUB),
    ("• AI가 자동으로 쿼리 생성 및 실행", False, TEXT_SUB),
    ("• 마크다운 표로 즉시 결과 제공", False, TEXT_SUB),
    ("• 응답 시간: 10~40초", False, TEXT_SUB),
    ("", False, TEXT_SUB),
    ("→ 누구나 즉시 데이터 조회 가능", True, GREEN),
    ("→ IT 의존도 대폭 감소", True, GREEN),
], 7.1, 1.9, 5.4, 4.0, size=15)

# 중간 화살표
text(sl, "→", 6.3, 3.6, 0.7, 0.6, size=30, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 3 — 시나리오 (영업팀 사용 예시)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "실제 사용 시나리오", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

# 페르소나
box(sl, 0.4, 1.15, 12.5, 0.85, fill=RGBColor(0x16,0x18,0x28), line_color=ACCENT, line_width=Pt(1))
text(sl, "👤  영업팀장 김과장  —  오늘 아침 경영진 보고 전, CRM 데이터를 빠르게 확인해야 합니다.", 0.65, 1.25, 12.0, 0.55, size=14, color=TEXT_SUB)

# 질문 말풍선 (사용자)
box(sl, 0.4, 2.2, 7.5, 0.75, fill=RGBColor(0x1E,0x20,0x3A), line_color=ACCENT, line_width=Pt(1))
text(sl, '💬  "이번 달 수주 현황 상위 5개 거래처 알려줘"', 0.65, 2.32, 7.0, 0.45, size=14, bold=True, color=ACCENT2)

# 처리 단계
for i, (icon, label) in enumerate([
    ("🔍", "질문 분석"),
    ("🔄", "쿼리 생성"),
    ("🗄️", "DB 조회"),
    ("📊", "결과 정리"),
]):
    bx = 0.4 + i * 3.1
    box(sl, bx, 3.2, 2.8, 0.85, fill=SURFACE, line_color=ACCENT2, line_width=Pt(1))
    text(sl, f"{icon}  {label}", bx+0.15, 3.32, 2.5, 0.5, size=13, color=TEXT_SUB)
    if i < 3:
        text(sl, "→", bx+2.85, 3.48, 0.35, 0.35, size=18, bold=True, color=ACCENT)

# 응답 말풍선 (AI)
box(sl, 0.4, 4.25, 12.5, 2.7, fill=RGBColor(0x05,0x14,0x10), line_color=GREEN, line_width=Pt(1.5))
text(sl, "🤖  AI 응답", 0.65, 4.35, 3, 0.4, size=13, bold=True, color=GREEN)
multiline(sl, [
    ("이번 달 수주 상위 5개 거래처입니다.", False, TEXT),
    ("", False, TEXT),
    ("  순위    거래처명             수주금액         수주건수", True, ACCENT2),
    ("  ────────────────────────────────────────────────", False, TEXT_SUB),
    ("  1위     (주)삼성전자         185,000,000원    3건", False, TEXT),
    ("  2위     현대자동차(주)       142,500,000원    2건", False, TEXT),
    ("  3위     LG화학               98,700,000원     4건", False, TEXT),
], 0.65, 4.75, 11.8, 2.1, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 4 — 아키텍처 전체 구성도
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "시스템 아키텍처", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

# 레이어 헤더
layers = [
    (0.3,  "🌐 프론트엔드",   "React SPA",             ACCENT),
    (3.1,  "⚙️  백엔드 서버",  "Node.js + Express",     RGBColor(0xA7,0x8B,0xFA)),
    (5.9,  "🤖 Claude Code", "AI 엔진 (CLI)",           GREEN),
    (8.7,  "🔌 MCP 레이어",   "Dataverse MCP",          YELLOW),
    (11.5, "🗄️ Azure",        "Dataverse CRM DB",       RED),
]

for bx, title, sub, clr in layers:
    box(sl, bx, 1.2, 2.5, 5.6, fill=SURFACE, line_color=clr, line_width=Pt(2))
    # 상단 색상 바
    box(sl, bx, 1.2, 2.5, 0.35, fill=clr)
    text(sl, title, bx+0.1, 1.22, 2.3, 0.3, size=11, bold=True, color=BG_DARK)
    text(sl, sub,   bx+0.15, 1.65, 2.2, 0.35, size=11, bold=True, color=clr)

# 각 레이어 내용
contents = [
    [("채팅 UI", TEXT), ("SSE 수신", TEXT), ("결과 렌더링", TEXT), ("세션 관리", TEXT)],
    [("Rate Limit", TEXT), ("API 키 인증", GREEN), ("동시접속 제한", GREEN), ("쓰기 차단", RED)],
    [("CLAUDE.md 로드", TEXT), ("자연어 이해", TEXT), ("MCP 도구 선택", TEXT), ("답변 생성", TEXT)],
    [("read_query", TEXT), ("search", TEXT), ("describe", TEXT), ("file_download", TEXT)],
    [("CRM 거래처", TEXT), ("영업기회", TEXT), ("수주 데이터", TEXT), ("서비스 케이스", TEXT)],
]

for i, (bx, _, _, clr) in enumerate(layers):
    for j, (item, clr2) in enumerate(contents[i]):
        box(sl, bx+0.12, 2.05+j*0.82, 2.26, 0.65, fill=RGBColor(0x12,0x14,0x1E), line_color=None)
        text(sl, item, bx+0.22, 2.12+j*0.82, 2.1, 0.5, size=11, color=clr2)

# 화살표
for ax in [2.8, 5.6, 8.4, 11.2]:
    text(sl, "→", ax, 3.7, 0.4, 0.4, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# 하단 설명
box(sl, 0.3, 7.0, 12.7, 0.35, fill=RGBColor(0x12,0x14,0x1E))
text(sl, "핵심: 서버가 Claude Code CLI를 child process로 실행 → CLAUDE.md 지침에 따라 MCP 도구로 Azure Dataverse 조회 → SSE로 실시간 스트리밍",
     0.5, 7.03, 12.3, 0.28, size=11, color=TEXT_SUB)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 5 — MCP란 무엇인가
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "MCP (Model Context Protocol) 이란?", 0.5, 0.3, 12, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

text(sl, "AI 모델이 외부 도구·데이터·서비스를 표준화된 방식으로 사용할 수 있게 해주는 프로토콜",
     0.5, 1.15, 12.3, 0.55, size=16, color=TEXT_SUB)

# 비유
box(sl, 0.4, 1.85, 12.5, 1.1, fill=RGBColor(0x16,0x18,0x28), line_color=ACCENT, line_width=Pt(1))
text(sl, "💡 쉬운 비유", 0.65, 1.93, 2.5, 0.35, size=13, bold=True, color=ACCENT2)
text(sl, '"스마트폰 앱이 카메라·GPS·연락처를 표준 API로 사용하듯, AI가 DB·파일·API를 MCP로 표준 접근합니다."',
     0.65, 2.28, 12.0, 0.5, size=14, color=TEXT)

# MCP 도구 3가지
items = [
    ("mcp__dataverse__read_query", "OData/FetchXML 쿼리 실행\n실제 CRM 데이터 조회", ACCENT),
    ("mcp__dataverse__search", "전체 텍스트 검색\n키워드로 레코드 탐색", RGBColor(0xA7,0x8B,0xFA)),
    ("mcp__dataverse__describe", "테이블 스키마 조회\n컬럼명·타입 파악", GREEN),
]
for i, (name, desc, clr) in enumerate(items):
    bx = 0.4 + i * 4.25
    box(sl, bx, 3.2, 4.0, 2.0, fill=SURFACE, line_color=clr, line_width=Pt(1.5))
    text(sl, name, bx+0.2, 3.3, 3.7, 0.45, size=11, bold=True, color=clr)
    text(sl, desc, bx+0.2, 3.82, 3.7, 1.2, size=13, color=TEXT_SUB)

# 보안
box(sl, 0.4, 5.45, 12.5, 1.75, fill=RGBColor(0x1A,0x0A,0x0A), line_color=RED, line_width=Pt(1.5))
text(sl, "🔒  보안 규칙 — 읽기 전용 (조회만 허용)", 0.65, 5.55, 8, 0.4, size=15, bold=True, color=RED)
multiline(sl, [
    ("✅  허용:  read_query  /  search  /  search_data  /  describe  /  file_download", False, GREEN),
    ("⛔  차단:  create_record  /  update_record  /  delete_record  /  create_table  /  delete_table  (서버 코드에서 강제 차단)", False, RED),
], 0.65, 6.0, 12.0, 1.1, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 6 — 요청 처리 흐름 (7단계)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "요청 처리 흐름", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

steps = [
    (ACCENT,  "1  사용자 질문",      "브라우저에서 POST /api/chat 전송\nsessionId + message"),
    (RGBColor(0x7C,0x3A,0xED), "2  서버 검증",      "Rate Limit · API Key 인증 · 동시접속 확인\n초과 시 429 즉시 반환"),
    (RGBColor(0x06,0xB6,0xD4), "3  Claude 실행",    "claude CLI spawn\nCLAUDE.md 지침 자동 로드"),
    (GREEN,   "4  MCP 도구 호출",    "read_query 등 실행\nOData로 Azure Dataverse 조회"),
    (YELLOW,  "5  보안 차단 검사",   "쓰기 도구 감지 시\n즉시 프로세스 kill + 에러 반환"),
    (RGBColor(0xF4,0x72,0x73), "6  답변 생성",       "데이터 기반 자연어 답변\n마크다운 표 형식"),
    (RGBColor(0xA7,0x8B,0xFA), "7  SSE 스트리밍",   "text 이벤트 실시간 전송\n글자 단위 렌더링"),
]

for i, (clr, title, desc) in enumerate(steps):
    row, col = divmod(i, 4)
    bx = 0.35 + col * 3.25
    by = 1.25 + row * 2.8
    box(sl, bx, by, 3.0, 2.4, fill=SURFACE, line_color=clr, line_width=Pt(2))
    # 번호
    circle = sl.shapes.add_shape(9, Inches(bx+0.1), Inches(by+0.1), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(i+1); run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE
    text(sl, title, bx+0.7, by+0.12, 2.2, 0.4, size=13, bold=True, color=clr)
    text(sl, desc,  bx+0.15, by+0.65, 2.75, 1.6, size=12, color=TEXT_SUB)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 7 — 보안 구조
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "보안 구조 — 3중 방어선", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

defenses = [
    ("🔑  1차 방어: API 키 인증",
     "API_KEY 환경변수 설정 시 활성화\n모든 /api/* 요청에 X-API-Key 헤더 필수\n미인증 요청 → 401 즉시 차단",
     ACCENT),
    ("⚡  2차 방어: 서버 코드 차단",
     "11개 쓰기 MCP 도구 이름을 Set으로 관리\ntool_use 이벤트 수신 시 실시간 감지\n감지 즉시 Claude 프로세스 kill + 에러 SSE",
     RED),
    ("📋  3차 방어: CLAUDE.md 지침",
     "Claude 자체 행동 지침 파일\n쓰기 금지 도구 명시\n사용자 요청이 와도 거절 + 이유 설명",
     YELLOW),
]

for i, (title, desc, clr) in enumerate(defenses):
    by = 1.3 + i * 1.9
    box(sl, 0.4, by, 12.5, 1.65, fill=SURFACE, line_color=clr, line_width=Pt(2))
    box(sl, 0.4, by, 0.12, 1.65, fill=clr)
    text(sl, title, 0.65, by+0.12, 6.5, 0.45, size=15, bold=True, color=clr)
    text(sl, desc,  0.65, by+0.62, 11.8, 0.95, size=13, color=TEXT_SUB)

# 추가 — Rate Limit
box(sl, 0.4, 7.0, 12.5, 0.35, fill=RGBColor(0x12,0x14,0x1E))
text(sl, "+ Rate Limit: 분당 IP당 20 요청 제한  /  동시 접속: 최대 5 프로세스 동시 실행  /  세션 TTL: 24시간 후 자동 만료",
     0.6, 7.03, 12.1, 0.28, size=11, color=TEXT_SUB)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 8 — 핵심 코드 설명
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "핵심 코드 구조", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

# 좌측 - Claude spawn
box(sl, 0.4, 1.2, 6.0, 2.6, fill=RGBColor(0x10,0x12,0x1E), line_color=ACCENT, line_width=Pt(1.5))
text(sl, "Claude Code 실행 방식", 0.6, 1.28, 5.6, 0.38, size=14, bold=True, color=ACCENT2)
code1 = """spawn('claude', [
  '-p', message,          // 사용자 질문
  '--output-format',
  'stream-json',          // 실시간 JSON 스트림
  '--resume', sessionId,  // 대화 이어받기
  '--dangerously-skip-
   permissions'           // MCP 자동 허용
], { cwd: projectRoot })  // CLAUDE.md 자동 로드"""
text(sl, code1, 0.55, 1.72, 5.7, 2.0, size=10.5, color=RGBColor(0xA5,0xB4,0xFC))

# 우측 - 쓰기 차단
box(sl, 6.9, 1.2, 6.0, 2.6, fill=RGBColor(0x14,0x08,0x08), line_color=RED, line_width=Pt(1.5))
text(sl, "쓰기 도구 실시간 차단", 7.1, 1.28, 5.6, 0.38, size=14, bold=True, color=RED)
code2 = """if (block.type === 'tool_use') {
  if (WRITE_TOOLS.has(block.name)) {
    // 즉시 프로세스 종료
    claude.kill()
    send({ type: 'error',
      message: '⛔ 변경 차단됨' })
    return
  }
  // 허용된 도구만 계속 실행
}"""
text(sl, code2, 7.05, 1.72, 5.7, 2.0, size=10.5, color=RGBColor(0xFC,0xA5,0xA5))

# 하단 - SSE 스트리밍
box(sl, 0.4, 4.05, 6.0, 1.85, fill=RGBColor(0x05,0x14,0x10), line_color=GREEN, line_width=Pt(1.5))
text(sl, "SSE 실시간 스트리밍", 0.6, 4.13, 5.6, 0.38, size=14, bold=True, color=GREEN)
code3 = """// 글자가 추가될 때마다 델타 전송
if (text.length > lastText.length) {
  send({ type: 'text',
    text: text.slice(lastText.length) })
  lastText = text
}"""
text(sl, code3, 0.55, 4.57, 5.7, 1.2, size=10.5, color=RGBColor(0x6E,0xE7,0xB7))

# 하단 - 세마포어
box(sl, 6.9, 4.05, 6.0, 1.85, fill=RGBColor(0x14,0x12,0x04), line_color=YELLOW, line_width=Pt(1.5))
text(sl, "동시 접속 세마포어 (최대 5개)", 7.1, 4.13, 5.6, 0.38, size=14, bold=True, color=YELLOW)
code4 = """// 초과 시 큐에서 대기
await acquireSemaphore()
const claude = spawn(...)
claude.on('close', () => {
  releaseSemaphore()  // 다음 요청 해제
})"""
text(sl, code4, 7.05, 4.57, 5.7, 1.2, size=10.5, color=RGBColor(0xFD,0xE6,0x8A))

# 파일 구조 요약
box(sl, 0.4, 6.1, 12.5, 1.25, fill=RGBColor(0x12,0x14,0x1E))
text(sl, "주요 파일:  server/index.ts (서버 핵심)  ·  CLAUDE.md (AI 지침)  ·  src/components/ (React UI)  ·  shared/types.ts (공통 타입)  ·  data/schema.json (스키마 캐시)",
     0.6, 6.18, 12.1, 0.45, size=12, color=TEXT_SUB)
text(sl, "설정 파일:  .env (환경변수)  ·  ecosystem.config.js (PM2)  ·  scripts/nginx.conf.example (nginx)",
     0.6, 6.65, 12.1, 0.35, size=12, color=TEXT_SUB)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 9 — 운영 구성
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "운영 구성", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

# PM2
box(sl, 0.4, 1.2, 3.9, 2.8, fill=SURFACE, line_color=GREEN, line_width=Pt(1.5))
text(sl, "🚀  PM2 프로세스 관리", 0.6, 1.28, 3.6, 0.4, size=13, bold=True, color=GREEN)
multiline(sl, [
    ("pm2 start ecosystem.config.js", True, RGBColor(0x6E,0xE7,0xB7)),
    ("  → 서버 기동", False, TEXT_SUB),
    ("pm2 restart crm-mcp", True, RGBColor(0x6E,0xE7,0xB7)),
    ("  → 재시작", False, TEXT_SUB),
    ("pm2 logs crm-mcp", True, RGBColor(0x6E,0xE7,0xB7)),
    ("  → 실시간 로그", False, TEXT_SUB),
    ("pm2 monit", True, RGBColor(0x6E,0xE7,0xB7)),
    ("  → CPU/메모리 모니터", False, TEXT_SUB),
], 0.6, 1.75, 3.65, 2.1, size=12)

# 환경변수
box(sl, 4.6, 1.2, 4.3, 2.8, fill=SURFACE, line_color=ACCENT2, line_width=Pt(1.5))
text(sl, "⚙️  주요 환경변수 (.env)", 4.8, 1.28, 4.0, 0.4, size=13, bold=True, color=ACCENT2)
multiline(sl, [
    ("PORT=3000", True, RGBColor(0xA5,0xB4,0xFC)),
    ("CHAT_TIMEOUT_MS=120000", True, RGBColor(0xA5,0xB4,0xFC)),
    ("RATE_LIMIT_MAX=20", True, RGBColor(0xA5,0xB4,0xFC)),
    ("MAX_CONCURRENT_CLAUDE=5", True, RGBColor(0x6E,0xE7,0xB7)),
    ("MAX_SESSIONS=200", True, RGBColor(0x6E,0xE7,0xB7)),
    ("API_KEY=****", True, YELLOW),
], 4.8, 1.75, 4.0, 2.1, size=12)

# 로그
box(sl, 9.2, 1.2, 3.9, 2.8, fill=SURFACE, line_color=YELLOW, line_width=Pt(1.5))
text(sl, "📋  로그 관리", 9.4, 1.28, 3.6, 0.4, size=13, bold=True, color=YELLOW)
multiline(sl, [
    ("logs/app.log", True, RGBColor(0xFD,0xE6,0x8A)),
    ("  → 전체 앱 로그 (JSON)", False, TEXT_SUB),
    ("logs/error.log", True, RGBColor(0xFD,0xE6,0x8A)),
    ("  → 에러 전용", False, TEXT_SUB),
    ("30일 자동 보관 후 삭제", False, TEXT_SUB),
    ("/api/logs → 최근 100건 조회", False, TEXT_SUB),
], 9.4, 1.75, 3.65, 2.1, size=12)

# 스키마 캐시
box(sl, 0.4, 4.25, 12.5, 1.55, fill=RGBColor(0x16,0x18,0x28), line_color=ACCENT, line_width=Pt(1))
text(sl, "📦  스키마 캐시 (data/schema.json)", 0.65, 4.33, 6, 0.4, size=14, bold=True, color=ACCENT2)
multiline(sl, [
    ("• 테이블 describe 결과를 파일에 저장 → 서버 재시작 후에도 즉시 응답 가능",  False, TEXT_SUB),
    ("• 배포 시 warmup_schema.sh 실행 권장 (단, Azure Dataverse 연결 상태에서만)",  False, TEXT_SUB),
], 0.65, 4.78, 12.1, 0.95, size=13)

# nginx
box(sl, 0.4, 6.0, 12.5, 1.35, fill=RGBColor(0x12,0x14,0x1E), line_color=TEXT_SUB, line_width=Pt(1))
text(sl, "🌐  nginx 리버스 프록시 (scripts/nginx.conf.example)", 0.65, 6.08, 9, 0.4, size=14, bold=True, color=TEXT_SUB)
text(sl, "SSE 스트리밍을 위해 proxy_buffering off  /  proxy_read_timeout 330s  /  HTTPS 전환 시 TLS 설정 주석 해제",
     0.65, 6.52, 12.1, 0.35, size=12, color=TEXT_SUB)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 10 — 확장 가능성
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
text(sl, "확장 가능성", 0.5, 0.3, 10, 0.6, size=28, bold=True, color=WHITE)
divider(sl, 1.0)

text(sl, "현재 구조는 Azure Dataverse 전용이 아닙니다 — MCP 서버만 교체하면 어느 DB든 연결 가능",
     0.5, 1.1, 12.3, 0.5, size=15, color=TEXT_SUB)

dbs = [
    ("Azure Dataverse", "Dataverse MCP\n(현재 운영 중)", GREEN, "✅"),
    ("PostgreSQL\n/ MySQL", "fastmcp 서버\n(코드 준비됨)", ACCENT2, "⚙️"),
    ("사내 ERP\n/ 사내 DB", "커스텀 MCP 서버\n(개발 필요)", YELLOW, "🔧"),
    ("REST API\n/ 외부 서비스", "HTTP MCP 서버\n(표준 프로토콜)", RGBColor(0xF4,0x72,0x73), "🌐"),
]

for i, (db, method, clr, icon) in enumerate(dbs):
    bx = 0.4 + i * 3.2
    box(sl, bx, 1.85, 3.0, 3.2, fill=SURFACE, line_color=clr, line_width=Pt(2))
    text(sl, icon, bx+1.2, 1.95, 0.7, 0.55, size=22, align=PP_ALIGN.CENTER, color=clr)
    text(sl, db, bx+0.15, 2.6, 2.75, 0.65, size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
    text(sl, method, bx+0.15, 3.4, 2.75, 0.85, size=12, color=TEXT_SUB, align=PP_ALIGN.CENTER)

# 교체 포인트
box(sl, 0.4, 5.25, 12.5, 2.1, fill=RGBColor(0x16,0x18,0x28), line_color=ACCENT, line_width=Pt(1))
text(sl, "DB 교체 시 변경 필요 파일 (4개)", 0.65, 5.33, 7, 0.4, size=14, bold=True, color=ACCENT2)
multiline(sl, [
    ("①  CLAUDE.md  — MCP 도구명 변경, 테이블 목록 재작성",       False, TEXT_SUB),
    ("②  server/index.ts  — WRITE_TOOLS Set 도구 이름 업데이트",  False, TEXT_SUB),
    ("③  data/schema.json  — 스키마 캐시 초기화 후 재적재",        False, TEXT_SUB),
    ("④  ~/.claude/ (MCP 설정)  — MCP 서버 연결 설정 변경",        False, TEXT_SUB),
], 0.65, 5.78, 12.1, 1.5, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 11 — 정리
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
box(sl, 0, 0, 0.12, 7.5, fill=ACCENT)

text(sl, "핵심 요약", 0.5, 0.5, 10, 0.6, size=32, bold=True, color=WHITE)
divider(sl, 1.25, x=0.5)

summaries = [
    ("🤖",  "Cloud LLM 활용",    "Anthropic Claude API를 서버에서 직접 호출\n별도 GPU 서버 불필요, 운영 비용 최소화"),
    ("🔌",  "MCP로 DB 연결",     "Azure Dataverse를 MCP 도구로 연결\nSQL 없이 자연어로 CRM 데이터 조회"),
    ("🔒",  "3중 보안",          "API 키 인증 + 서버 쓰기 차단 + CLAUDE.md 지침\nAzure 데이터는 절대 변경 불가"),
    ("⚡",  "실시간 스트리밍",   "SSE로 답변을 글자 단위 실시간 전송\n사용자 체감 응답속도 향상"),
    ("🔧",  "DB 확장 용이",      "MCP 서버만 교체하면 어느 DB든 연결\nPostgreSQL·사내 ERP 모두 지원 가능"),
    ("📊",  "운영 안정성",       "PM2 프로세스 관리 + 동시접속 제한\n로그 자동 관리 + Rate Limit 보호"),
]

for i, (icon, title, desc) in enumerate(summaries):
    row, col = divmod(i, 3)
    bx = 0.4 + col * 4.25
    by = 1.5 + row * 2.5
    box(sl, bx, by, 4.0, 2.25, fill=SURFACE, line_color=ACCENT, line_width=Pt(1))
    text(sl, f"{icon}  {title}", bx+0.2, by+0.15, 3.65, 0.45, size=14, bold=True, color=ACCENT2)
    text(sl, desc, bx+0.2, by+0.68, 3.65, 1.45, size=12, color=TEXT_SUB)

# 저장
out_path = os.path.join(os.path.dirname(__file__), '..', 'mcp_azure.pptx')
prs.save(out_path)
print(f"저장 완료: {os.path.abspath(out_path)}")
print(f"슬라이드 수: {len(prs.slides)}")
